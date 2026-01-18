# =========================================================
# SOCIAL MEDIA REPORTING AUTOMATION – FINAL STREAMLIT APP
# =========================================================

import streamlit as st
import pandas as pd
import numpy as np

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ---------------------------------------------------------
# STREAMLIT UI
# ---------------------------------------------------------

st.set_page_config(page_title="Social Media Reporting", layout="wide")
st.title("📊 Social Media Reporting Automation")

uploaded_file = st.file_uploader(
    "Upload your social media Excel file (.xlsx)",
    type=["xlsx"]
)

if uploaded_file is None:
    st.info("Please upload an Excel file to begin.")
    st.stop()

# ---------------------------------------------------------
# CONFIG
# ---------------------------------------------------------

PRIORITY_BRANDS = [
    "Similac",
    "Ensure",
    "Pediasure",
    "Pedialyte",
    "Juven",
    "Glucerna"
]

EXCEL_OUTPUT = "Final_Social_Media_Report.xlsx"
PPT_OUTPUT = "Final_Social_Media_Trends.pptx"

# ---------------------------------------------------------
# LOAD DATA
# ---------------------------------------------------------

df = pd.read_excel(uploaded_file)

# ---------------------------------------------------------
# STANDARDIZE COLUMNS
# ---------------------------------------------------------

df = df.rename(columns={
    "Social Channel/Platform": "platform",
    "Brand Name/Category Name": "brand",
    "Type of Post (Branded, Influencer,Creators, Organic, Shop)": "post_type",
    "Post Format (Video, Reels, Shorts, Images, Carousels)": "format",
    "Video Plays": "views",
    "Followers": "followers",
    "Influencer Tier": "influencer_tier"
})

for c in ["Likes", "Comments", "views", "followers"]:
    df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)

df["engagement"] = df["Likes"] + df["Comments"]

df["brand"] = df["brand"].astype(str).str.strip().str.title()

df["post_type"] = (
    df["post_type"]
    .astype(str)
    .str.title()
    .replace({
        "Tagged": "Influencer",
        "Creators": "Creator"
    })
)

df["format"] = (
    df["format"]
    .astype(str)
    .str.title()
    .replace({
        "Reels": "Shorts",
        "Images": "Image",
        "Sidecar": "Carousel",
        "Image+Video": "Carousel"
    })
)

df["is_dynamic"] = df["format"].isin(["Video", "Shorts"])
df["is_paid"] = df["post_type"].isin(["Branded", "Influencer", "Creator", "Shop"])

# ---------------------------------------------------------
# BRAND ORDERING
# ---------------------------------------------------------

def sort_brands(df, volume_col):
    df = df.copy()
    df["__priority"] = df["brand"].apply(
        lambda x: PRIORITY_BRANDS.index(x)
        if x in PRIORITY_BRANDS else len(PRIORITY_BRANDS)
    )
    df = df.sort_values(
        by=["__priority", volume_col],
        ascending=[True, False]
    )
    return df.drop(columns="__priority")

# ---------------------------------------------------------
# METRIC FUNCTIONS
# ---------------------------------------------------------

def overall_metrics(x):
    video_eng = x.loc[x["is_dynamic"], "engagement"].sum()
    video_views = x.loc[x["is_dynamic"], "views"].sum()

    return pd.Series({
        "Total Posts": len(x),
        "Total Engagement": x["engagement"].sum(),
        "Video Engagement": video_eng,
        "Total Views": video_views,
        "Avg ER%": video_eng / video_views if video_views > 0 else None
    })

def format_metrics(x):
    fmt = x["format"].iloc[0]

    if fmt in ["Video", "Shorts"]:
        eng = x["engagement"].sum()
        views = x["views"].sum()
        er = eng / views if views > 0 else None
    else:
        eng = x["engagement"].sum()
        foll = x["followers"].sum()
        er = eng / foll if foll > 0 else None

    return pd.Series({
        "Post Count": len(x),
        "Avg ER%": er
    })

# ---------------------------------------------------------
# CONTENT SCOPES
# ---------------------------------------------------------

scopes = {
    "Paid": df[df["is_paid"]],
    "Organic": df[df["post_type"] == "Organic"],
    "Branded": df[df["post_type"] == "Branded"],
    "Influencer": df[df["post_type"] == "Influencer"],
    "Creator": df[df["post_type"] == "Creator"]
}

outputs = {}

# ---------------------------------------------------------
# GENERATE TABLES
# ---------------------------------------------------------

for scope, sdf in scopes.items():

    # ---------- OVERALL ----------
    overall = sdf.groupby("brand").apply(overall_metrics).reset_index()
    outputs[f"{scope}_Brand_Overall"] = sort_brands(overall, "Total Posts")

    # ---------- FORMAT (LONG) ----------
    format_long = (
        sdf.groupby(["brand", "format"])
        .apply(format_metrics)
        .reset_index()
    )

    # ---------- FORMAT – POST COUNT (WIDE) ----------
    format_post_count = (
        format_long
        .pivot_table(
            index="brand",
            columns="format",
            values="Post Count",
            aggfunc="sum",
            fill_value=0
        )
        .reset_index()
    )

    # ---------- FORMAT – AVG ER% (WIDE) ----------
    format_avg_er = (
        format_long
        .pivot_table(
            index="brand",
            columns="format",
            values="Avg ER%",
            aggfunc="first"
        )
        .reset_index()
        .fillna("-")
    )

    outputs[f"{scope}_Format_PostCount"] = sort_brands(format_post_count, "Video" if "Video" in format_post_count.columns else format_post_count.columns[1])
    outputs[f"{scope}_Format_AvgER"] = sort_brands(format_avg_er, format_avg_er.columns[1])

    # ---------- TYPE ----------
    typ = sdf.groupby(["brand", "post_type"]).apply(overall_metrics).reset_index()
    outputs[f"{scope}_Brand_Type"] = sort_brands(typ, "Total Posts")

    # ---------- SOURCE ----------
    src = sdf.groupby(["brand", "platform"]).apply(overall_metrics).reset_index()
    outputs[f"{scope}_Brand_Source"] = sort_brands(src, "Total Posts")

    # ---------- WEEKLY ----------
    wk = sdf.groupby(["brand", "Week"]).apply(overall_metrics).reset_index()
    outputs[f"{scope}_Brand_Weekly"] = wk

# ---------------------------------------------------------
# INFLUENCER TIER ANALYSIS
# ---------------------------------------------------------

tier = (
    df[df["post_type"] == "Influencer"]
    .groupby(["brand", "influencer_tier"])
    .apply(overall_metrics)
    .reset_index()
)

outputs["Influencer_Tier_Analysis"] = sort_brands(tier, "Total Posts")

# ---------------------------------------------------------
# EXPORT EXCEL
# ---------------------------------------------------------

with pd.ExcelWriter(EXCEL_OUTPUT, engine="openpyxl") as writer:
    for name, table in outputs.items():
        table.to_excel(writer, sheet_name=name[:31], index=False)

# ---------------------------------------------------------
# PPT – WEEK ON WEEK TRENDLINES
# ---------------------------------------------------------

prs = Presentation()

def add_trend_slide(df, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    chart_data = CategoryChartData()
    weeks = sorted(df["Week"].dropna().unique())
    chart_data.categories = weeks

    for brand in PRIORITY_BRANDS:
        bdf = df[df["brand"] == brand].sort_values("Week")
        if not bdf.empty:
            chart_data.add_series(
                brand,
                bdf["Avg ER%"].fillna(0).tolist()
            )

    slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        x=0,
        y=0,
        cx=9144000,
        cy=4572000,
        chart_data=chart_data
    )

add_trend_slide(outputs["Paid_Brand_Weekly"], "Paid Content – WoW Avg ER%")
add_trend_slide(outputs["Branded_Brand_Weekly"], "Branded Content – WoW Avg ER%")
add_trend_slide(outputs["Influencer_Brand_Weekly"], "Influencer Content – WoW Avg ER%")
add_trend_slide(outputs["Creator_Brand_Weekly"], "Creator Content – WoW Avg ER%")

prs.save(PPT_OUTPUT)

# ---------------------------------------------------------
# DOWNLOADS
# ---------------------------------------------------------

st.success("Report generated successfully!")

with open(EXCEL_OUTPUT, "rb") as f:
    st.download_button("📥 Download Excel Report", f, file_name=EXCEL_OUTPUT)

with open(PPT_OUTPUT, "rb") as f:
    st.download_button("📥 Download PPT (WoW Trends)", f, file_name=PPT_OUTPUT)